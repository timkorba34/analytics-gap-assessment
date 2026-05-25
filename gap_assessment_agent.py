import json
import os
import streamlit as st
from openai import OpenAI
import PyPDF2
import pandas as pd
import io

from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from tavily import TavilyClient
from graphviz import Digraph


# --------------------
# Page Setup
# --------------------
st.set_page_config(page_title="AI Gap Assessment Builder", layout="wide")

import base64

def get_base64(file_path):
    with open(file_path, "rb") as f:
        return base64.b64encode(f.read()).decode()

bg_image = get_base64("copilot_bg.png")
robot_image = get_base64("robot_bg.png")

# ==========================
# App Styling
# ==========================

st.markdown(
    f"""
<style>

/* ==========================
   APP THEME
========================== */

.stApp {{
    background: linear-gradient(
        rgba(5,10,25,.96),
        rgba(5,10,25,.98)
    );
    background-attachment: fixed;
    color:white;
}}

label,
.stMarkdown p,
.stMarkdown span,
label[data-testid="stWidgetLabel"] {{
    color:#ffffff !important;
}}

label[data-testid="stWidgetLabel"] {{
    font-weight:600 !important;
}}


/* ==========================
   HERO SECTION
========================== */

.hero-container {{
    display:flex;
    justify-content:space-between;
    align-items:center;

    background:linear-gradient(
        90deg,
        rgba(10,15,40,.95),
        rgba(15,23,42,.82)
    );

    padding:50px;
    border-radius:28px;

    border:1px solid rgba(255,255,255,.08);

    margin-bottom:30px;
}}

.hero-left {{
    width:58%;
}}

.hero-right {{
    width:38%;
    text-align:right;
}}

.hero-right img {{
    width:100%;
    max-width:420px;
}}

.hero-title {{
    font-size:64px;
    font-weight:800;
    color:white;
    line-height:1.1;
}}

.hero-subtitle {{
    font-size:22px;
    color:#c7d2fe;
    line-height:1.6;
}}


/* ==========================
   GLASS CONTAINERS
========================== */

.glass {{
    background:rgba(15,23,42,.72);
    backdrop-filter:blur(12px);

    padding:30px;
    border-radius:20px;

    border:1px solid rgba(255,255,255,.12);

    box-shadow:0 0 30px rgba(0,0,0,.4);

    margin-bottom:25px;
}}


/* ==========================
   FEATURE CARDS
========================== */

.feature-card {{
    background:rgba(20,30,60,.72);

    padding:24px;
    border-radius:18px;

    border:1px solid rgba(255,255,255,.12);

    min-height:160px;
}}

.feature-card h3 {{
    color:white !important;
    font-size:26px;
}}

.feature-card p {{
    color:#e5e7eb !important;
    font-size:17px;
}}


/* ==========================
   INPUTS
========================== */

.stTextInput input,
.stTextArea textarea {{
    background-color:rgba(15,23,42,.85)!important;

    color:white!important;

    border:1px solid rgba(255,255,255,.25)!important;

    border-radius:10px!important;
}}

.stSelectbox div[data-baseweb="select"] > div {{
    background-color:rgba(15,23,42,.85)!important;

    color:white!important;

    border:1px solid rgba(255,255,255,.25)!important;

    border-radius:10px!important;
}}


/* ==========================
   FILE UPLOADER
========================== */

[data-testid="stFileUploader"] section {{
    background-color:rgba(10,20,45,.90)!important;

    border:1px dashed rgba(255,255,255,.35)!important;

    border-radius:14px!important;
}}

[data-testid="stFileUploader"] section:hover {{
    border:1px solid #4da6ff!important;
}}

[data-testid="stFileUploader"] section + div {{
    background:transparent!important;
}}

[data-testid="stFileUploader"] ul,
[data-testid="stFileUploader"] li,
[data-testid="stFileUploader"] li > div,
[data-testid="stFileUploader"] li div {{

    background-color:rgba(17,26,51,.98)!important;

    color:white!important;

    border-radius:12px!important;
}}

[data-testid="stFileUploader"] li {{
    border:1px solid rgba(96,165,250,.65)!important;
    margin-bottom:8px!important;
}}

[data-testid="stFileUploader"] span,
[data-testid="stFileUploader"] p,
[data-testid="stFileUploader"] small,
[data-testid="stFileUploader"] svg {{

    color:white!important;
    fill:white!important;
}}

[data-testid="stFileUploader"] button {{

    background:#0b1f44!important;

    color:white!important;

    border:1px solid rgba(96,165,250,.65)!important;
}}


/* ==========================
   BUTTONS
========================== */

div.stButton > button,
[data-testid="stDownloadButton"] button {{

    background:linear-gradient(
        90deg,
        #2563eb,
        #7c3aed
    ) !important;

    color:white!important;

    border:none!important;

    border-radius:14px!important;

    height:55px!important;

    font-size:18px!important;

    font-weight:600!important;

    width:100%;
}}

div.stButton > button:hover {{
    background:linear-gradient(
        90deg,
        #1d4ed8,
        #6d28d9
    )!important;
}}

button * {{
    color:white!important;
}}

</style>
""",
unsafe_allow_html=True
)


# ==========================
# HERO
# ==========================

st.markdown(
f"""
<div class="hero-container">

<div class="hero-left">

<div class="hero-title">
Analytics Modernization Assessment Copilot
</div>

<div class="hero-subtitle">
AI-powered executive analytics assessments,
modernization roadmaps, S/4HANA readiness analysis,
and actionable remediation planning.
</div>

</div>

<div class="hero-right">
<img src="data:image/png;base64,{robot_image}">
</div>

</div>
""",
unsafe_allow_html=True
)

# --------------------
# OpenAI Setup
# --------------------
api_key = st.secrets.get("OPENAI_API_KEY", None) or os.getenv("OPENAI_API_KEY")

if not api_key:
    st.error("OPENAI_API_KEY is missing. Add it in Streamlit Cloud Secrets.")
    st.stop()

client = OpenAI(api_key=api_key)

#----------------------
# Tavili Setup
#______________________



tavily_api_key = st.secrets.get("TAVILY_API_KEY", None) or os.getenv("TAVILY_API_KEY")

tavily_client = None
if tavily_api_key:
    tavily_client = TavilyClient(api_key=tavily_api_key)


# --------------------
# Initialize Session State
# --------------------
defaults = {
    "assessment_data": None,
    "word_doc": None,
    "ppt_file": None,
    "email_text": None,
}

for key, value in defaults.items():
    if key not in st.session_state:
        st.session_state[key] = value

# --------------------
# Company Search - Drop Down
# --------------------
def search_companies(query):
    try:
        response = tavily_client.search(
            query=f"{query} official company name official website",
            search_depth="basic",
            max_results=5,
            exclude_domains=[
                "linkedin.com",
                "facebook.com",
                "twitter.com",
                "x.com",
                "glassdoor.com",
                "zoominfo.com",
                "seamless.ai",
                "rocketreach.co",
                "wikipedia.org"
            ]
        )

        search_context = ""

        for result in response.get("results", []):
            search_context += f"""
Title: {result.get("title", "")}
URL: {result.get("url", "")}
Content: {result.get("content", "")}
"""

        prompt = f"""
You are identifying the correct official company name.

User typed:
{query}

Search results:
{search_context}

Return only the official company name.
Do not return page titles.
Do not return descriptions.
Do not return URLs.
Do not return LinkedIn, ZoomInfo, Seamless, employee names, or marketing text.
Return one clean company name only.

Example:
User typed: PIM Brands
Return: PIM Brands, Inc.
"""

        result = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Return only one clean official company name."},
                {"role": "user", "content": prompt}
            ],
            temperature=0
        )

        company_name = result.choices[0].message.content.strip()

        return [company_name]

    except Exception:
        return [query.strip()]
        
# --------------------
# UI Inputs
# --------------------
st.markdown("<div class='glass'>", unsafe_allow_html=True)

client_search = st.text_input("Search Client")

company_options = []

if client_search:
    company_options = search_companies(client_search)

if company_options:

    client_name = st.selectbox(
        "Select Company",
        company_options
    )

else:
    client_name = client_search
    
industry = st.selectbox(
    "Industry",
    [
        "Select an industry",
        "Accounting",
        "Airlines and Aviation",
        "Alternative Medicine",
        "Animation and Post-Production",
        "Apparel Manufacturing",
        "Architecture and Planning",
        "Automotive",
        "Banking",
        "Biotechnology Research",
        "Broadcast Media Production and Distribution",
        "Building Materials",
        "Business Consulting and Services",
        "Chemical Manufacturing",
        "Civil Engineering",
        "Computer and Network Security",
        "Construction",
        "Consumer Services",
        "Data Infrastructure and Analytics",
        "Defense and Space Manufacturing",
        "Education",
        "Education Administration Programs",
        "Electric Power Generation",
        "Entertainment Providers",
        "Environmental Services",
        "Events Services",
        "Executive Offices",
        "Facilities Services",
        "Farming",
        "Financial Services",
        "Food and Beverage Manufacturing",
        "Government Administration",
        "Healthcare",
        "Hospitals and Health Care",
        "Hospitality",
        "Human Resources Services",
        "Industrial Machinery Manufacturing",
        "Information Services",
        "Information Technology & Services",
        "Insurance",
        "Investment Banking",
        "IT Services and IT Consulting",
        "Law Practice",
        "Legal Services",
        "Life Sciences",
        "Logistics and Supply Chain",
        "Machinery Manufacturing",
        "Manufacturing",
        "Market Research",
        "Medical Equipment Manufacturing",
        "Mining",
        "Non-profit Organizations",
        "Oil and Gas",
        "Pharmaceutical Manufacturing",
        "Primary and Secondary Education",
        "Professional Services",
        "Public Relations and Communications Services",
        "Real Estate",
        "Research Services",
        "Retail",
        "Retail Apparel and Fashion",
        "Semiconductor Manufacturing",
        "Software Development",
        "Sports",
        "Telecommunications",
        "Transportation/Trucking/Railroad",
        "Travel Arrangements",
        "Utilities",
        "Warehousing and Storage",
        "Wholesale",
        "Wireless Services",
        "Other"
    ]
)

source_system = st.selectbox(
    "Source System",
    [
        "SAP ECC",
        "SAP S/4HANA",
        "Oracle EBS",
        "Oracle Fusion",
        "JD Edwards",
        "Microsoft Dynamics",
        "Infor",
        "NetSuite",
        "Multiple ERP Systems",
        "Other"
    ]
)

migration_type = st.selectbox(
    "Migration Type",
    [
        "Brownfield",
        "Greenfield",
        "Smartfield / Selective Data Transition",
        "Not Yet Defined"
    ]
)

target_source_system = st.selectbox(
    "Target Source System",
    [
        "SAP S/4HANA",
        "Oracle Fusion",
        "SAP ECC (No Migration)",
        "Microsoft Dynamics",
        "Multiple ERP Systems",
        "Other"
    ]
)

assessment_type = st.selectbox(
    "Assessment Type",
    [
        "Analytics Gap Assessment",
        "Analytics Modernization Assessment",
        "S/4HANA Impact Assessment",
        "AI Readiness Assessment",
        "Data Governance Assessment",
        "Reporting Rationalization Assessment"
    ]
)

uploaded_files = st.file_uploader(
    "Upload Discovery Notes / Supporting Files",
    type=["txt", "csv", "pdf", "xls", "xlsx", "doc", "docx"],
    accept_multiple_files=True
)



notes = st.text_area("Paste Additional Notes", height=250)

safe_client_name = client_name.strip().replace(" ", "_") if client_name else "Client"

st.markdown("</div>", unsafe_allow_html=True)

# --------------------
# File Reader
# --------------------
def read_uploaded_files(files):
    content = ""

    if not files:
        return content

    for file in files:
        content += f"\n\n--- FILE: {file.name} ---\n"

        file_name = file.name.lower()

        try:
            if file_name.endswith(".txt") or file.type == "text/plain":
                content += file.read().decode("utf-8", errors="ignore")

            elif file_name.endswith(".csv") or file.type == "text/csv":
                df = pd.read_csv(file)
                content += df.head(25).to_string(index=False)

            elif file_name.endswith(".xlsx"):
                excel_file = pd.ExcelFile(file, engine="openpyxl")
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    content += f"\n\n--- SHEET: {sheet_name} ---\n"
                    content += df.head(20).to_string(index=False)

            elif file_name.endswith(".xls"):
                excel_file = pd.ExcelFile(file, engine="xlrd")
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    content += f"\n\n--- SHEET: {sheet_name} ---\n"
                    content += df.head(20).to_string(index=False)

            elif file_name.endswith(".pdf") or file.type == "application/pdf":
                reader = PyPDF2.PdfReader(file)
                for i, page in enumerate(reader.pages, start=1):
                    content += f"\n\n--- PAGE {i} ---\n"
                    content += page.extract_text() or ""

            elif file_name.endswith(".docx"):

                doc = Document(file)

                for paragraph in doc.paragraphs:
                    content += paragraph.text + "\n"

            else:
                content += f"\nUnsupported file type: {file.type}"

        except Exception as e:
            content += f"\nError reading file: {str(e)}"

    return content

# --------------------
# OpenAI Retry Helper
# --------------------
from openai import RateLimitError, APIError, APITimeoutError
import time
import json

response_format={"type": "json_object"}

def call_openai_with_retry(messages, model="gpt-4o-mini"):
    for attempt in range(3):
        try:
            return client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=0.2,
                max_tokens=4000,
                response_format={"type": "json_object"}
            )

        except RateLimitError:
            if attempt < 2:
                time.sleep(2 ** attempt)
            else:
                st.error("OpenAI rate limit reached. Check billing/quota or reduce upload size.")
                return None

        except (APIError, APITimeoutError) as e:
            st.error(f"OpenAI API error: {str(e)}")
            return None

# --------------------
# Research Company Information
# --------------------

def research_company(company_name, industry):
    if not tavily_client or not company_name:
        return ""

    query = f"{company_name} company overview industry products revenue locations acquisitions strategy {industry}"

    try:
        results = tavily_client.search(
            query=query,
            search_depth="basic",
            max_results=5
        )

        research_text = ""

        for item in results.get("results", []):
            title = item.get("title", "")
            url = item.get("url", "")
            content = item.get("content", "")

            research_text += f"\nTitle: {title}\nURL: {url}\nSummary: {content}\n"

        return research_text[:6000]

    except Exception as e:
        return f"Company research unavailable: {str(e)}"

# --------------------
# Shared Prompt Builder
# --------------------
def build_base_context(client_name, industry, assessment_type, notes, file_content, company_research):
    return f"""
CLIENT INFORMATION
Client Name: {client_name}
Industry: {industry}
Assessment Type: {assessment_type}
Source System: {source_system}
Migration Type: {migration_type}
Target Source System: {target_source_system}
Assessment Type: {assessment_type}

Adjust recommendations, impacted reports, dependencies,
table mappings, architecture, and roadmap based on these selections.

Transformation Context:
Current Source System: {source_system}
Target Source System: {target_source_system}
Migration Type: {migration_type}

The assessment must explicitly reference this transformation path throughout the report. 
Do not generically say “S/4HANA migration.” Explain how the selected migration type changes

PUBLIC COMPANY RESEARCH
{company_research}

DISCOVERY NOTES
{notes[:4000]}

SUPPORTING FILE CONTENT
{file_content[:12000]}

{COMPANY_RESEARCH_REQUIREMENTS}

{TABLE_SUMMARY_REQUIREMENTS}



WRITING RULES
You are a senior consulting partner creating a paid executive deliverable.

Write with a direct, commercial, boardroom-ready point of view.

Before writing, infer the company's operating model from the discovery notes, uploaded files, research, and industry.

Do not write generic statements that could apply to any company.

Every issue must explain:
- where it shows up in the business
- what is happening operationally
- who is impacted
- what decision is delayed, wrong, or harder to make
- why it matters financially or operationally

No placeholders.
No "To be validated."
No empty strings.
No empty arrays.
No markdown.
No code fences.

Return valid JSON only.

All tables must be arrays of flat row objects.
Do not create nested dictionaries inside table cells.
"""

# --------------------
# Recommend Embedded Reports
# --------------------

def recommend_embedded_reports(report_inventory):

    recommendations=[]

    for report in report_inventory:

        prompt=f"""
        Current Report:
        {report}

        Search SAP standard embedded analytics,
        SAP Fiori analytical applications,
        CDS analytical queries,
        and recommend:

        - Best replacement
        - Confidence score
        - Retain/Replace/Rebuild/Retire
        - Reason
        """

        result = client.chat.completions.create(...)

        recommendations.append(result)

    return recommendations

# --------------------
# Assessment Context
# --------------------
ASSESSMENT_CONTEXT = {
    "Analytics Gap Assessment": """
You are creating a premium executive-level Analytics Gap Assessment.

The customer is currently undergoing or planning an SAP S/4HANA transformation initiative. The purpose is to evaluate the current-state data, analytics, reporting, governance, ownership, and source-system environment to identify gaps, risks, and business disruption areas caused by source and reporting changes.

Focus on:
- current-state reporting risks
- S/4HANA source-system impact
- report breakage risk
- data model and KPI disruption
- historical reporting challenges
- governance and ownership gaps
- business decision-making risk
""",

    "Analytics Modernization Roadmap": """
You are creating a premium executive-level Analytics Modernization Roadmap.

The customer is working toward an enterprise analytics modernization journey. The purpose is to define the future-state analytics vision, platform strategy, governance model, modernization priorities, roadmap, and business value.

Focus on:
- future-state analytics architecture
- platform consolidation
- reporting modernization
- governance maturity
- self-service analytics
- cloud/data platform strategy
- phased execution roadmap
- business value realization
""",

    "AI Opportunity Assessment": """
You are creating a premium executive-level AI Opportunity Assessment.

The customer is looking to understand AI readiness, practical AI opportunities, and where AI can create value across operations, analytics, automation, and decision-making.

Focus on:
- AI readiness
- data maturity
- AI governance
- automation opportunities
- decision-support use cases
- practical GenAI opportunities
- implementation complexity
- phased AI roadmap
- business value realization
"""
}

COMPANY_RESEARCH_REQUIREMENTS = """
Use the public company research, discovery notes, uploaded files, and industry context to make the assessment company-specific.

Explain the customer's operating model, likely business complexity, reporting needs, ERP/data landscape, and transformation drivers.

Avoid generic language. Do not fabricate specific counts, revenue, locations, systems, reports, or years of history unless provided.
"""

TABLE_SUMMARY_REQUIREMENTS = """
After every table, generate a 1-2 paragraph executive narrative explaining what the table means to the customer.

The narrative must explain:
- why leadership should care
- operational implications
- business risk
- reporting or analytics impact
- financial or decision-making impact
- recommended action

Do not simply repeat the table. Interpret the findings.
"""

# --------------------
# Generate One Section
# --------------------
def generate_section_json(
    client_name,
    industry,
    assessment_type,
    notes,
    file_content,
    company_research,
    section_config,
    assessment_context
):
    base_context = build_base_context(
        client_name,
        industry,
        assessment_type,
        notes,
        file_content,
        company_research
    )

    cross_section_rules = "\n".join(
    f"- {rule}"
    for rule in SECTION_INSTRUCTIONS.get("cross_section_rules", [])
)

    required_keys = "\n".join(section_config["keys"])

    prompt = f"""
{assessment_context}



{COMPANY_RESEARCH_REQUIREMENTS}

{TABLE_SUMMARY_REQUIREMENTS}

{base_context}

GLOBAL CROSS-SECTION RULES:
{cross_section_rules}

SECTION TO GENERATE:
{section_config["section_name"]}

SECTION INSTRUCTIONS:
{section_config["instructions"]}

REQUIRED JSON KEYS:
{required_keys}

If the required keys contain "transformation_context", generate it as an array.

Format:

"transformation_context": [
  {{
    "area": "ERP",
    "current_state": "SAP ECC",
    "target_state": "SAP S/4HANA",
    "migration_impact": "Brownfield migration retains existing configurations and customizations, requiring validation of existing reporting dependencies and custom code."
  }},
  {{
    "area": "Finance",
    "current_state": "BKPF/BSEG/FAGLFLEXA",
    "target_state": "ACDOCA",
    "migration_impact": "Universal Journal consolidation impacts financial reporting logic and report structures."
  }},
  {{
    "area": "Inventory",
    "current_state": "MKPF/MSEG",
    "target_state": "MATDOC",
    "migration_impact": "Inventory reporting logic and historical movement analysis require validation."
  }},
  {{
    "area": "Sales",
    "current_state": "VBAK/VBAP/VBRK/VBRP",
    "target_state": "S/4HANA CDS Views",
    "migration_impact": "Order-to-cash reporting dependencies and KPIs require revalidation."
  }},
  {{
    "area": "Reporting",
    "current_state": "BusinessObjects / Excel / BW",
    "target_state": "S/4HANA Embedded Analytics / Datasphere",
    "migration_impact": "Reports should be categorized into retain, remediate, rebuild, or retire."
  }}
]

Return only these keys in one valid JSON object.

Requirements:
- Every required key must be populated
- No placeholder text
- No generic consulting language
- Use company-specific business context
- Tie findings to operational and reporting impacts
- Explain business risks and executive implications
- Explain S/4HANA impacts where relevant
- Interpret findings like an executive consultant
"""
    
    messages = [
        {
            "role": "system",
            "content": "Return one valid JSON object only. No markdown. No commentary."
        },
        {
            "role": "user",
            "content": prompt
        }
    ]

    response = call_openai_with_retry(messages)

    if response is None:
        return {}

    raw = response.choices[0].message.content.strip()

    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        st.error(f"Invalid JSON returned for section: {section_config['section_name']}")
        st.code(raw[:4000])
        return {}


# --------------------
# Generate Full Assessment - Multiple Calls
# --------------------
def generate_assessment_json(
    client_name,
    industry,
    assessment_type,
    notes,
    file_content,
    company_research
):
    framework = ASSESSMENT_FRAMEWORKS.get(assessment_type)
    assessment_context = ASSESSMENT_CONTEXT.get(assessment_type, "")

    if not framework:
        st.error(f"No framework found for assessment type: {assessment_type}")
        return {}

    assessment_context = ASSESSMENT_CONTEXT.get(assessment_type, "")

    final_data = {}

    for section in framework["sections"]:
        with st.spinner(f"Generating {section['section_name']}..."):
            section_data = generate_section_json(
                client_name,
                industry,
                assessment_type,
                notes,
                file_content,
                company_research,
                section,
                assessment_context
            )

        final_data.update(section_data)

    return final_data

# --------------------
# Table of Contents Helper
# --------------------
def set_update_fields_on_open(doc):
    settings = doc.settings.element
    update_fields = OxmlElement("w:updateFields")
    update_fields.set(qn("w:val"), "true")
    settings.append(update_fields)
    
# --------------------
# Word Table Formatting Helpers
# --------------------

TABLE_COUNTER = {"count": 0}

def set_cell_shading(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_margins(cell, top=80, start=80, bottom=80, end=80):
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    tc_mar = tc_pr.first_child_found_in("w:tcMar")

    if tc_mar is None:
        tc_mar = OxmlElement("w:tcMar")
        tc_pr.append(tc_mar)

    for m, v in {
        "top": top,
        "start": start,
        "bottom": bottom,
        "end": end,
    }.items():
        node = tc_mar.find(qn(f"w:{m}"))
        if node is None:
            node = OxmlElement(f"w:{m}")
            tc_mar.append(node)

        node.set(qn("w:w"), str(v))
        node.set(qn("w:type"), "dxa")


def format_table(table):
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = True

    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.TOP
            set_cell_margins(cell)

            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT

                for run in paragraph.runs:
                    run.font.size = Pt(8)

            if row_idx == 0:
                set_cell_shading(cell, "D9D9D9")

                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        run.font.size = Pt(8)


def add_table_caption(doc):
    TABLE_COUNTER["count"] += 1

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    run = p.add_run(f"Table {TABLE_COUNTER['count']}")
    run.bold = True
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(80, 80, 80)


def add_table_of_contents(doc):
    doc.add_heading("Table of Contents", level=1)

    paragraph = doc.add_paragraph()
    run = paragraph.add_run()

    fld_char1 = OxmlElement("w:fldChar")
    fld_char1.set(qn("w:fldCharType"), "begin")

    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = 'TOC \\o "1-3" \\h \\z \\u'

    fld_char2 = OxmlElement("w:fldChar")
    fld_char2.set(qn("w:fldCharType"), "separate")

    fld_char3 = OxmlElement("w:fldChar")
    fld_char3.set(qn("w:fldCharType"), "end")

    run._r.append(fld_char1)
    run._r.append(instr_text)
    run._r.append(fld_char2)
    run._r.append(fld_char3)

    doc.add_page_break()

# --------------------
# Word Helpers
# --------------------
def add_heading(doc, text, level=1):
    doc.add_heading(text, level=level)


def add_paragraph(doc, text):
    if not text:
        return

    if isinstance(text, dict):
        text = json.dumps(text, indent=2)
    elif isinstance(text, list):
        text = "\n".join([str(item) for item in text])
    else:
        text = str(text)

    doc.add_paragraph(text)


def add_table_from_records(doc, records):
    if not records:
        return

    if isinstance(records, str):
        doc.add_paragraph(records)
        return

    if isinstance(records, dict):
        records = [records]

    if not isinstance(records, list) or len(records) == 0:
        doc.add_paragraph("No records were generated for this section.")
        return

    if isinstance(records[0], str):
        for item in records:
            doc.add_paragraph(str(item), style="List Bullet")
        return

    if not isinstance(records[0], dict):
        doc.add_paragraph(str(records))
        return

    headers = list(records[0].keys())

    table = doc.add_table(rows=1, cols=len(headers))

    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = str(h)

    for record in records:
        row = table.add_row().cells

        for i, h in enumerate(headers):
            value = record.get(h, "")

            if isinstance(value, dict):
                value = json.dumps(value, indent=2)
            elif isinstance(value, list):
                value = ", ".join([str(x) for x in value])
            else:
                value = str(value)

            row[i].text = value

    # IMPORTANT: these must be OUTSIDE both loops
    format_table(table)
    add_table_caption(doc)


# --------------------
# Build Word Document
# --------------------
def build_docx(data, client_name, assessment_type):
    TABLE_COUNTER["count"] = 0

    doc = Document()

    doc.add_heading(f"{client_name or 'Client'} {assessment_type}", 0)
    add_table_of_contents(doc)

    # Common executive front-end
    add_heading(doc, "1. Engagement Overview", 1)
    doc.add_paragraph("")
    add_paragraph(doc, data.get("engagement_overview_text", ""))

    add_heading(doc, "2. Executive Summary", 1)
    doc.add_paragraph("")
    add_paragraph(doc, data.get("executive_summary_text", ""))

    if data.get("top_priorities"):
        add_heading(doc, "Executive Priorities", 2)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("top_priorities", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("top_priorities_text", ""))  

    set_update_fields_on_open(doc)

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)

    return output

    # --------------------
    # Analytics Gap Assessment
    # --------------------
    if assessment_type == "Analytics Gap Assessment":

        add_heading(doc, "3. S/4HANA Transformation Context", 1)
        doc.add_paragraph("")
        add_table_from_records(
        doc,
        data.get("transformation_context", []))
        doc.add_paragraph("")
        add_paragraph(
        doc,
        data.get("transformation_context_summary", ""))

        add_heading(doc, "4. Current Analytics Ecosystem", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("current_system_inventory", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("current_data_flow_summary", ""))

        add_heading(doc, "5. Reporting Dependency Map", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("reporting_dependency_map", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("architecture_risk_summary", ""))

        add_heading(doc, "6. Analytics Complexity and Operational Risk", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("analytics_complexity_snapshot", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("analytics_complexity_text", ""))

        add_heading(doc, "7. Gap Severity Heatmap", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("gap_severity_heatmap", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("gap_observations_text", ""))

        add_heading(doc, "8. Reporting Inventory Summary", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("reporting_landscape_summary", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("reporting_inventory_text", ""))

        add_heading(doc, "9. Report Modernization and Embedded Analytics Recommendation", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("report_replacement_matrix", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("report_replacement_text", ""))

        add_heading(doc, "10. S/4HANA Reporting Impact", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("s4_impact_summary", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("s4_reporting_impact_text", ""))

        add_heading(doc, "11. Gap Analysis Summary", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("gap_analysis_summary", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("key_gaps_text", ""))

        add_heading(doc, "12. Opportunity Areas", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("improvement_opportunity_summary", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("opportunity_areas_text", ""))

        add_heading(doc, "13. Business Value", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("potential_impact_summary", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("business_value_text", ""))

        add_heading(doc, "14. Gap Remediation Roadmap", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("s4_analytics_roadmap", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("s4_analytics_roadmap_text", ""))

        add_heading(doc, "15. Recommended Remediation Actions and Execution Plan", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("recommended_focus_areas", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("recommended_next_steps_text", ""))

        add_heading(doc, "16. Appendix A — Reporting Inventory", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("appendix_reporting_inventory", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("appendix_reporting_inventory_text", ""))

        add_heading(doc, "17. Appendix B — S/4 Reporting Impact Analysis", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("appendix_s4_impact_analysis", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("appendix_s4_impact_analysis_text", ""))

        add_heading(doc, "18. Appendix C — Reporting Overlap Analysis", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("appendix_reporting_overlap_analysis", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("appendix_reporting_overlap_analysis_text", ""))

        add_heading(doc, "19. Appendix D — Data Source Mapping", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("appendix_data_source_mapping", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("appendix_data_source_mapping_text", ""))

        add_heading(doc, "20. Appendix E — Critical Reports", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("appendix_critical_reports", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("appendix_critical_reports_text", ""))

        add_heading(doc, "Critical Report Summary", 2)
        doc.add_paragraph("")
        add_paragraph(doc, data.get("critical_report_summary", ""))

        add_heading(doc, "21. Appendix F — Analytics Stakeholder Map", 1)
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("analytics_responsibility_model", []))
        add_table_from_records(doc, data.get("stakeholder_interview_summary", []))
        doc.add_paragraph("")
        add_table_from_records(doc, data.get("responsibility_gaps", []))
        doc.add_paragraph("")
        add_paragraph(doc, data.get("analytics_ownership_overview_text", ""))

    # --------------------
    # Analytics Modernization Roadmap
    # --------------------
    elif assessment_type == "Analytics Modernization Roadmap":

        add_heading(doc, "3. Modernization Drivers", 1)
        add_table_from_records(doc, data.get("modernization_drivers", []))
        doc.add_paragraph("")

        add_heading(doc, "4. Current-State Architecture", 1)
        add_table_from_records(doc, data.get("current_state_architecture", []))
        doc.add_paragraph("")

        add_heading(doc, "5. Future-State Architecture", 1)
        add_table_from_records(doc, data.get("future_state_architecture", []))
        doc.add_paragraph("")

        add_heading(doc, "6. Capability Gap Summary", 1)
        add_table_from_records(doc, data.get("capability_gap_summary", []))
        doc.add_paragraph("")

        add_heading(doc, "7. Platform Recommendations", 1)
        add_table_from_records(doc, data.get("platform_recommendations", []))
        doc.add_paragraph("")

        add_heading(doc, "8. Workstream Plan", 1)
        add_table_from_records(doc, data.get("workstream_plan", []))
        doc.add_paragraph("")

        add_heading(doc, "9. Risk Mitigation Plan", 1)
        add_table_from_records(doc, data.get("risk_mitigation_plan", []))
        doc.add_paragraph("")

        add_heading(doc, "10. Investment Summary", 1)
        add_table_from_records(doc, data.get("investment_summary", []))
        doc.add_paragraph("")

        add_heading(doc, "11. Business Value", 1)
        add_paragraph(doc, data.get("business_value_text", ""))
        add_table_from_records(doc, data.get("potential_impact_summary", []))
        doc.add_paragraph("")

    # --------------------
    # AI Opportunity Assessment
    # --------------------
    elif assessment_type == "AI Opportunity Assessment":

        add_heading(doc, "3. Top AI Opportunities", 1)
        add_table_from_records(doc, data.get("top_ai_opportunities", []))
        doc.add_paragraph("")

        add_heading(doc, "4. AI Use Case Inventory", 1)
        add_table_from_records(doc, data.get("ai_use_case_inventory", []))
        doc.add_paragraph("")

        add_heading(doc, "5. Automation Candidates", 1)
        add_table_from_records(doc, data.get("automation_candidates", []))
        doc.add_paragraph("")

        add_heading(doc, "6. Decision Support Opportunities", 1)
        add_table_from_records(doc, data.get("decision_support_opportunities", []))
        doc.add_paragraph("")

        add_heading(doc, "7. Data Readiness Summary", 1)
        add_table_from_records(doc, data.get("data_readiness_summary", []))
        doc.add_paragraph("")

        add_heading(doc, "8. AI Roadmap", 1)
        add_table_from_records(doc, data.get("ai_roadmap", []))
        doc.add_paragraph("")

        add_heading(doc, "9. Risk and Governance Considerations", 1)
        add_table_from_records(doc, data.get("risk_and_governance_considerations", []))
        doc.add_paragraph("")

        add_heading(doc, "10. Business Value", 1)
        add_paragraph(doc, data.get("business_value_text", ""))
        add_table_from_records(doc, data.get("potential_impact_summary", []))
        doc.add_paragraph("")

        add_heading(doc, "11. Recommended Next Steps", 1)
        add_paragraph(doc, data.get("recommended_next_steps_text", ""))
        doc.add_paragraph("")

        add_heading(doc, "Key Observations", 1)
        add_paragraph(doc, data.get("key_observations_text", ""))
        doc.add_paragraph("")

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)

    return output

def build_ppt(data, client_name):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 Title
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"{client_name} Analytics Gap Assessment"

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = data.get("executive_summary_text", "Executive summary unavailable.")
    p.font.size = Pt(18)

    # Slide 2 Key Gaps
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Key Analytics Gaps"

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
    tf = tx.text_frame

    gaps = data.get("gap_analysis_summary", [])

    if isinstance(gaps, str):
        gaps = [{"Gap": gaps, "Business Impact": ""}]
    elif isinstance(gaps, dict):
        gaps = [gaps]
    elif not isinstance(gaps, list):
        gaps = []

    for gap in gaps[:6]:
        p = tf.add_paragraph()

        if isinstance(gap, dict):
            p.text = f"• {gap.get('Gap', 'Gap')} – {gap.get('Business Impact', '')}"
        else:
            p.text = f"• {str(gap)}"

        p.font.size = Pt(16)

    # Slide 3 Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Recommended Next Steps"

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
    tf = tx.text_frame

    focus = data.get("recommended_focus_areas", [])

    if isinstance(focus, str):
        focus = [{"Focus Area": focus, "Recommended Next Step": ""}]
    elif isinstance(focus, dict):
        focus = [focus]
    elif not isinstance(focus, list):
        focus = []

    for item in focus[:6]:
        p = tf.add_paragraph()

        if isinstance(item, dict):
            p.text = f"• {item.get('Focus Area', '')} – {item.get('Recommended Next Step', '')}"
        else:
            p.text = f"• {str(item)}"

        p.font.size = Pt(16)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return output
    
def build_exec_email(data, client_name):
    summary = data.get("executive_summary_text", "")
    
    email = f"""
Subject: {client_name} Analytics Gap Assessment – Executive Summary

Team,

We completed the initial analytics gap assessment for {client_name}.

Key observations:
{summary}

Top priorities identified:
1. Centralize reporting and KPI visibility
2. Improve data integration across systems
3. Enable forecasting and operational analytics
4. Build scalable analytics foundation for growth

Recommended next step:
Conduct a focused strategy workshop and roadmap session to prioritize quick wins and transformation initiatives.

Regards,
Consulting Team
"""
    return email

# --------------------
# Company Research Default
# --------------------
company_research = ""

# --------------------
# Output Validation
# --------------------
def validate_output(data, assessment_type):
    if not data:
        st.error("Validation failed: data object is empty.")
        return False

    framework = ASSESSMENT_FRAMEWORKS.get(assessment_type)

    if not framework:
        st.error(f"Validation failed: no framework found for {assessment_type}")
        return False

    required_keys = []

    for section in framework["sections"]:
        required_keys.extend(section["keys"])

    for key in required_keys:
        value = data.get(key)

        if value is None:
            st.error(f"Missing key: {key}")
            return False

        if isinstance(value, str):
            if not value.strip():
                st.error(f"Empty string key: {key}")
                return False
            if "to be validated" in value.lower():
                st.error(f"Placeholder found in key: {key}")
                return False

        if isinstance(value, list):
            if len(value) == 0:
                st.error(f"Empty list key: {key}")
                return False
            if "to be validated" in str(value).lower():
                st.error(f"Placeholder found in list key: {key}")
                st.write(value)
                return False

        if isinstance(value, dict):
            if len(value) == 0:
                st.error(f"Empty dict key: {key}")
                return False
            if "to be validated" in str(value).lower():
                st.error(f"Placeholder found in dict key: {key}")
                return False

    return True

# --------------------
# Section Instructions by Document Header
# --------------------

SECTION_INSTRUCTIONS = {
    "Engagement Overview": """
engagement_overview_text:
Must be 2-3 substantial executive paragraphs.
Explain why the assessment was performed, the client business context, current source system, target source system, migration type, and why analytics/reporting risk matters.
""",

    "cross_section_rules": [
        "Do not restate previously identified gaps",
        "Reference earlier findings instead",
        "Company-specific examples must appear in every section",
        "Recommendations must explain capability before technology",
        "Avoid repeating same business risks more than once"
    ],

    "Executive Summary": """
executive_summary_text:
Must be 3-4 substantial executive paragraphs.
Summarize the highest-priority findings, business risk, reporting risk, governance gaps, S/4HANA impact, and urgency for remediation.

top_priorities:
Must be a table array with exactly 5 rows.
Columns: Priority, Why It Matters, Business Impact, Time Horizon, Executive Owner

top_priorities_text:
Must be 2 executive paragraphs interpreting the priority table.
""",

    "S/4HANA Transformation Context": """
transformation_context must be a table array.
Columns: Area, Current State, Target State, Migration Impact

transformation_context_summary:
Must be 1-2 executive paragraphs explaining source-to-target migration path, Brownfield/Greenfield/Smartfield impact, report breakage risk, source table impacts, business process impacts, and remediation considerations.
""",

    "Current Analytics Ecosystem": """
current_system_inventory must be a table array.
Columns: Business Area, Current Source System, Reporting Tool, Data Owner, Integration Method, Refresh Frequency, Key Dependency, Current Pain Point, S/4HANA Risk

current_data_flow_summary:
Must be 1-2 executive paragraphs explaining data flow, fragmentation, manual effort, reporting bottlenecks, operational risks, and S/4HANA disruption concerns.
""",

    "Reporting Dependency Map": """
reporting_dependency_map must be a table array.
Columns: Report / Dashboard, Primary Source System, Dependent Systems, Business Function, Criticality, Current Risk, S/4HANA Impact

architecture_risk_summary:
Must be 1-2 executive paragraphs explaining architecture weaknesses, integration concerns, reporting dependencies, scalability limitations, and business continuity risks.
""",

    "Analytics Complexity and Operational Risk": """
analytics_complexity_snapshot must be a table array.
Columns: Business Area, Data Sources, Reporting Tools, Complexity Level, Key Challenges, Operational Risk, Remediation Priority

analytics_complexity_text:
Must be 1-2 executive paragraphs explaining why the analytics environment is complex, where operational risk appears, and why leadership should address it before or during migration.
""",

    "Gap Severity Heatmap": """
gap_severity_heatmap must be a table array.
Columns: Gap Area, Severity Level, Impact, Business Risk, Recommended Priority

gap_observations_text:
Must be 1-2 executive paragraphs explaining the highest-risk gaps, why leadership should care, and what should be remediated first.
""",

    "Reporting Inventory Summary": """
reporting_landscape_summary must be a table array.
Columns: Report Type, Frequency, Users, Issues, S/4HANA Risk, Remediation Need

reporting_inventory_text:
Must be 1-2 executive paragraphs explaining the current reporting landscape, redundancy, manual effort, risk of report breakage, and need for inventory rationalization.
""",

    "Report Modernization and Embedded Analytics Recommendation": """
report_replacement_matrix must be a table array.

Columns:
Report, Tool, ABAP/T-Code, Area, Source Tables, Purpose, Fiori App, Embedded Query, Disposition, Target Option, Confidence, Rationale, Deep Dive?

Rules:
- Do not use "To be validated"
- Do not use "TBD"
- Do not leave cells blank
- If the exact value is unknown, use one of these approved values:
  - Unknown
  - Not Provided
  - Requires SME Review
  - Requires Technical Review
  - Not Applicable
- Source Tables must contain either known SAP tables/views or "Requires Technical Review"
- ABAP/T-Code must contain a known program/T-code or "Not Applicable"
- Fiori App must contain a likely app name or "Requires Functional Review"
- Embedded Query must contain a likely analytical query or "Requires Technical Review"

report_replacement_text:
Must be 1-2 executive paragraphs explaining which Excel, Power BI, ABAP, BW, BusinessObjects, or other reports should be retained, remediated, replaced, rebuilt, retired, or moved to modern analytics.
"""

    "S/4HANA Reporting Impact": """
s4_impact_summary must be a table array.
Columns: Impact Area, Current State, Target State, Risk, Reporting Impact, Remediation Action

s4_reporting_impact_text:
Must be 1-2 executive paragraphs explaining how the selected migration type affects reporting, table structures, KPIs, historical reporting, and business continuity.
""",

    "Gap Analysis Summary": """
gap_analysis_summary must be a table array.
Columns: Gap Area, Current State, Identified Risk, Recommended Action, Business Impact

key_gaps_text:
Must be 1-2 executive paragraphs summarizing the most important gaps and why they matter operationally, financially, and strategically.
""",

    "Opportunity Areas": """
improvement_opportunity_summary must be a table array.
Columns: Opportunity, Description, Business Value, Priority, Recommended Next Step

opportunity_areas_text:
Must be 1-2 executive paragraphs explaining how the client can use the transformation to modernize reporting, governance, integration, and analytics capabilities.
""",

    "Business Value": """
potential_impact_summary must be a table array.
Columns: Impact Area, Description, Business Value, Financial / Operational Impact, Priority

business_value_text:
Must be 1-2 executive paragraphs explaining the business value of remediation and modernization.
Do not fabricate financial metrics unless provided.
""",

    "Gap Remediation Roadmap": """
s4_analytics_roadmap must be a post-assessment gap remediation roadmap.
Do not include an Assessment phase. The assessment has already been completed.

Use phases such as:
1. Stabilize Critical Reporting Risk
2. Define Data Ownership and KPI Governance
3. Remediate S/4HANA Reporting Dependencies
4. Modernize Priority Reports and Data Flows
5. Operationalize Analytics Governance and Continuous Improvement

Columns:
Phase, Timeline, Gap Addressed, Remediation Actions, Expected Outcome, Business Value, Dependencies

s4_analytics_roadmap_text:
Must be 1-2 executive paragraphs explaining how the roadmap moves the client from assessment into execution.
""",

    "Recommended Remediation Actions and Execution Plan": """
recommended_focus_areas must be a highly actionable post-assessment execution plan tied directly to identified gaps.

Columns:
Recommendation Category, Gap Addressed, Recommended Action, Business Outcome, Priority, Suggested Owner, Execution Horizon, Potential Follow-On Deliverable

recommended_next_steps_text:
Must be 1-2 executive paragraphs explaining immediate actions, urgent risks, governance decisions, workstreams, funding alignment, and follow-on implementation activities.
""",

    "Appendix A — Reporting Inventory": """
appendix_reporting_inventory must be a table array.
Columns: Report Name, Report Type, Frequency, Owner, Data Source, Criticality, S/4HANA Risk

appendix_reporting_inventory_text:
Must be 1-2 paragraphs explaining why this inventory matters.
""",

    "Appendix B — S/4 Reporting Impact Analysis": """
appendix_s4_impact_analysis must be a table array.
Columns: Area, Current State, Target State, Impact, Risk Level, Remediation Consideration

appendix_s4_impact_analysis_text:
Must be 1-2 paragraphs explaining table, reporting, KPI, and source-to-target impact.
""",

    "Appendix C — Reporting Overlap Analysis": """
appendix_reporting_overlap_analysis must be a table array.
Columns: Report Name, Overlap With, Description, Business Impact, Recommended Action

appendix_reporting_overlap_analysis_text:
Must be 1-2 paragraphs explaining redundancy, confusion, and consolidation opportunities.
""",

    "Appendix D — Data Source Mapping": """
appendix_data_source_mapping must be a table array.
Columns: Report Name, Data Source, Dependency, Data Quality, S/4HANA Risk

appendix_data_source_mapping_text:
Must be 1-2 paragraphs explaining source dependencies and migration risk.
""",

    "Appendix E — Critical Reports": """
appendix_critical_reports must be a table array.
Columns: Report Name, Criticality Level, Impact on Decision, Impact of Failure, Remediation Priority

appendix_critical_reports_text:
Must be 1-2 paragraphs explaining why these reports must be prioritized.

critical_report_summary:
Must be 1 executive paragraph summarizing the critical report risk.
""",

    "Appendix F — Analytics Stakeholder Map": """
analytics_responsibility_model must be a table array.
Columns: Department, Report Name, Responsibility

stakeholder_interview_summary must be a table array.
Columns: Stakeholder, Role, Key Concerns

responsibility_gaps must be a table array.
Columns: Report Name, Gap, Business Impact, Recommended Action

analytics_ownership_overview_text:
Must be 1-2 paragraphs explaining ownership gaps and governance risk.

key_observations_text:
Must be 1-2 paragraphs summarizing appendix observations.
"""
}

# --------------------
# Assessment Frameworks
# --------------------

ASSESSMENT_FRAMEWORKS = {
    "Analytics Gap Assessment": {
        "title": "Analytics Gap Assessment",
        "sections": [
            {
                "section_name": "Engagement Overview",
                "keys": ["engagement_overview_text"],
                "instructions": SECTION_INSTRUCTIONS["Engagement Overview"]
            },
            {
                "section_name": "Executive Summary",
                "keys": ["executive_summary_text", "top_priorities", "top_priorities_text"],
                "instructions": SECTION_INSTRUCTIONS["Executive Summary"]
            },
            {
                "section_name": "S/4HANA Transformation Context",
                "keys": ["transformation_context", "transformation_context_summary"],
                "instructions": SECTION_INSTRUCTIONS["S/4HANA Transformation Context"]
            },
            {
                "section_name": "Current Analytics Ecosystem",
                "keys": ["current_system_inventory", "current_data_flow_summary"],
                "instructions": SECTION_INSTRUCTIONS["Current Analytics Ecosystem"]
            },
            {
                "section_name": "Reporting Dependency Map",
                "keys": ["reporting_dependency_map", "architecture_risk_summary"],
                "instructions": SECTION_INSTRUCTIONS["Reporting Dependency Map"]
            },
            {
                "section_name": "Analytics Complexity and Operational Risk",
                "keys": ["analytics_complexity_snapshot", "analytics_complexity_text"],
                "instructions": SECTION_INSTRUCTIONS["Analytics Complexity and Operational Risk"]
            },
            {
                "section_name": "Gap Severity Heatmap",
                "keys": ["gap_severity_heatmap", "gap_observations_text"],
                "instructions": SECTION_INSTRUCTIONS["Gap Severity Heatmap"]
            },
            {
                "section_name": "Reporting Inventory Summary",
                "keys": ["reporting_landscape_summary", "reporting_inventory_text"],
                "instructions": SECTION_INSTRUCTIONS["Reporting Inventory Summary"]
            },
            {
                "section_name": "Report Modernization and Embedded Analytics Recommendation",
                "keys": ["report_replacement_matrix", "report_replacement_text"],
                "instructions": SECTION_INSTRUCTIONS["Report Modernization and Embedded Analytics Recommendation"]
            },
            {
                "section_name": "S/4HANA Reporting Impact",
                "keys": ["s4_impact_summary", "s4_reporting_impact_text"],
                "instructions": SECTION_INSTRUCTIONS["S/4HANA Reporting Impact"]
            },
            {
                "section_name": "Gap Analysis Summary",
                "keys": ["gap_analysis_summary", "key_gaps_text"],
                "instructions": SECTION_INSTRUCTIONS["Gap Analysis Summary"]
            },
            {
                "section_name": "Opportunity Areas",
                "keys": ["improvement_opportunity_summary", "opportunity_areas_text"],
                "instructions": SECTION_INSTRUCTIONS["Opportunity Areas"]
            },
            {
                "section_name": "Business Value",
                "keys": ["potential_impact_summary", "business_value_text"],
                "instructions": SECTION_INSTRUCTIONS["Business Value"]
            },
            {
                "section_name": "Gap Remediation Roadmap",
                "keys": ["s4_analytics_roadmap", "s4_analytics_roadmap_text"],
                "instructions": SECTION_INSTRUCTIONS["Gap Remediation Roadmap"]
            },
            {
                "section_name": "Recommended Remediation Actions and Execution Plan",
                "keys": ["recommended_focus_areas", "recommended_next_steps_text"],
                "instructions": SECTION_INSTRUCTIONS["Recommended Remediation Actions and Execution Plan"]
            },
            {
                "section_name": "Appendix A — Reporting Inventory",
                "keys": ["appendix_reporting_inventory", "appendix_reporting_inventory_text"],
                "instructions": SECTION_INSTRUCTIONS["Appendix A — Reporting Inventory"]
            },
            {
                "section_name": "Appendix B — S/4 Reporting Impact Analysis",
                "keys": ["appendix_s4_impact_analysis", "appendix_s4_impact_analysis_text"],
                "instructions": SECTION_INSTRUCTIONS["Appendix B — S/4 Reporting Impact Analysis"]
            },
            {
                "section_name": "Appendix C — Reporting Overlap Analysis",
                "keys": ["appendix_reporting_overlap_analysis", "appendix_reporting_overlap_analysis_text"],
                "instructions": SECTION_INSTRUCTIONS["Appendix C — Reporting Overlap Analysis"]
            },
            {
                "section_name": "Appendix D — Data Source Mapping",
                "keys": ["appendix_data_source_mapping", "appendix_data_source_mapping_text"],
                "instructions": SECTION_INSTRUCTIONS["Appendix D — Data Source Mapping"]
            },
            {
                "section_name": "Appendix E — Critical Reports",
                "keys": ["appendix_critical_reports", "appendix_critical_reports_text", "critical_report_summary"],
                "instructions": SECTION_INSTRUCTIONS["Appendix E — Critical Reports"]
            },
            {
                "section_name": "Appendix F — Analytics Stakeholder Map",
                "keys": [
                    "analytics_responsibility_model",
                    "stakeholder_interview_summary",
                    "responsibility_gaps",
                    "analytics_ownership_overview_text",
                    "key_observations_text"
                ],
                "instructions": SECTION_INSTRUCTIONS["Appendix F — Analytics Stakeholder Map"]
            }
        ]
    }
}

# --------------------
# Generate Button
# --------------------
if st.button("Generate Assessment Outputs", key="main_generate_btn"):

    if not client_name:
        st.warning("Enter a client name first.")
    else:
        file_content = read_uploaded_files(uploaded_files)
        company_research = research_company(client_name, industry)

        with st.spinner("Generating assessment content..."):
            max_retries = 3
            data = None

            for attempt in range(max_retries + 1):
                data = generate_assessment_json(
                    client_name,
                    industry,
                    assessment_type,
                    notes,
                    file_content,
                    company_research
                )
        
                if validate_output(data, assessment_type):
                    break
                else:
                    st.warning(f"Regenerating output attempt {attempt + 1}: missing sections or placeholder text found.")
        
            if not validate_output(data, assessment_type):
                st.error("Failed to generate a complete assessment after retries.")
                data = {}

        st.session_state.assessment_data = data

        if data:
            with st.spinner("Creating Word document..."):
                st.session_state.word_doc = build_docx(data, client_name, assessment_type)

            with st.spinner("Creating PowerPoint deck..."):
                st.session_state.ppt_file = build_ppt(data, client_name)

            with st.spinner("Creating Executive Summary Email..."):
                st.session_state.email_text = build_exec_email(data, client_name)

            st.success("Assessment outputs generated successfully.")
        else:
            st.error("Assessment generation failed.")

# --------------------
# Download Buttons
# --------------------
if st.session_state.get("word_doc"):
    st.download_button(
        label="Download Word Document",
        data=st.session_state.word_doc.getvalue() if hasattr(st.session_state.word_doc, "getvalue") else st.session_state.word_doc,
        file_name=f"{safe_client_name}_Gap_Assessment_Report.docx",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        key="download_word_doc"
    )

if st.session_state.get("ppt_file"):
    st.download_button(
        label="Download PowerPoint Deck",
        data=st.session_state.ppt_file.getvalue() if hasattr(st.session_state.ppt_file, "getvalue") else st.session_state.ppt_file,
        file_name=f"{safe_client_name}_Gap_Assessment_Deck.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key="download_ppt_file"
    )

if st.session_state.get("email_text"):
    st.download_button(
        label="Download Executive Summary Email",
        data=st.session_state.email_text,
        file_name=f"{safe_client_name}_Executive_Summary_Email.txt",
        mime="text/plain",
        key="download_exec_email"
    )
